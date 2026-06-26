#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Genera la presentacion EP3 (Tienda de Perritos - ECS Fargate)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BLUE   = RGBColor(0x3B, 0x82, 0xF6)
DARK   = RGBColor(0x1F, 0x29, 0x37)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF3, 0xF4, 0xF6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def style_run(r, size, color=DARK, bold=False, italic=False, font="Calibri"):
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def header_band(slide, title, kicker=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    band.fill.solid(); band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    style_run(r, 28, WHITE, bold=True)
    if kicker:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = kicker
        style_run(r2, 13, RGBColor(0xDB, 0xEA, 0xFE))
    # numero de lamina
    return band


def bullets(slide, items, l=Inches(0.7), t=Inches(1.5), w=Inches(12), h=Inches(5.4), size=18):
    tb, tf = textbox(slide, l, t, w, h)
    first = True
    for txt, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(8)
        r = p.add_run()
        bullet = "•  " if lvl == 0 else "–  "
        r.text = bullet + txt
        style_run(r, size - lvl*2, DARK if lvl == 0 else GRAY, bold=(lvl == 0))
    return tb


def box(slide, l, t, w, h, text, fill, fg=WHITE, size=12, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = WHITE; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for i, line in enumerate(text.split("\n")):
        pp = p if i == 0 else tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = line
        style_run(r, size if i == 0 else size-2, fg, bold=(bold and i == 0))
    return sp


def arrow(slide, x1, y1, x2, y2, color=DARK):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    conn.line.color.rgb = color; conn.line.width = Pt(2.25)
    line_elem = conn.line._get_or_add_ln()
    tail = line_elem.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    line_elem.append(tail)
    return conn


# ───────────────────────── 1. PORTADA ─────────────────────────
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = BLUE; bg.line.fill.background(); bg.shadow.inherit = False
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.0), SW, Inches(2.5))
strip.fill.solid(); strip.fill.fore_color.rgb = WHITE; strip.line.fill.background(); strip.shadow.inherit = False
_, tf = textbox(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(3.0))
r = tf.paragraphs[0].add_run(); r.text = "Tienda de Perritos"
style_run(r, 54, WHITE, bold=True)
p = tf.add_paragraph(); r = p.add_run(); r.text = "Orquestacion y automatizacion en la nube con AWS ECS Fargate + CI/CD"
style_run(r, 24, RGBColor(0xDB, 0xEA, 0xFE))
p = tf.add_paragraph(); r = p.add_run(); r.text = "Evaluacion Parcial N3  -  Introduccion a Herramientas DevOps (ISY1101)"
style_run(r, 16, RGBColor(0xDB, 0xEA, 0xFE), italic=True)
_, tf = textbox(s, Inches(0.8), Inches(5.25), Inches(11.7), Inches(2.0))
for txt, sz, col, bd in [
    ("Caso: Innovatech Chile", 20, DARK, True),
    ("Integrantes: Gabriel Molina  -  Jefferson Rojas", 16, GRAY, False),
    ("Junio 2026  -  github.com/jeffplop/tienda-perritos", 14, GRAY, False),
]:
    p = tf.add_paragraph(); r = p.add_run(); r.text = txt; style_run(r, sz, col, bold=bd)
set_notes(s, "Saludar, presentarnos e indicar el tema: desplegamos una app de 3 capas como "
    "contenedores orquestados en AWS ECS Fargate, con un pipeline CI/CD automatizado en GitHub "
    "Actions. Recordar que el caso es la empresa Innovatech Chile que avanza hacia orquestacion "
    "productiva tras contenedorizar (EP2) y montar infraestructura base (EP1).")

# ───────────────────────── 2. AGENDA ─────────────────────────
s = add_slide(); header_band(s, "Agenda")
bullets(s, [
    ("Contexto del caso y objetivo", 0),
    ("Arquitectura general de la solucion", 0),
    ("Por que ECS Fargate (decisiones tecnicas)", 0),
    ("Configuracion del cluster: redes, security groups y roles IAM", 0),
    ("Despliegue de Frontend y Backend + comunicacion entre servicios", 0),
    ("Autoscaling", 0),
    ("Pipeline CI/CD: build -> push -> deploy (demo)", 0),
    ("Gestion de secrets, logs y metricas", 0),
    ("Problemas encontrados, conclusiones y proyeccion productiva", 0),
], size=20)
set_notes(s, "En 10-15 minutos vamos a recorrer la arquitectura, la configuracion del cluster, el "
    "despliegue, el autoscaling, el pipeline (con demo si el lab esta arriba) y cerramos con "
    "problemas/aprendizajes. Cada uno debe poder explicar cualquier parte: la defensa es individual.")

# ───────────────────────── 3. CONTEXTO ─────────────────────────
s = add_slide(); header_band(s, "Contexto del caso", "Innovatech Chile")
bullets(s, [
    ("EP1: infraestructura base en AWS.  EP2: contenedorizacion de la app.", 0),
    ("EP3: avanzar hacia orquestacion y automatizacion productiva.", 0),
    ("Objetivo: ejecutar la aplicacion de forma escalable, tolerante a fallos y automatizable.", 0),
    ("Automatizar por completo los despliegues desde GitHub.", 0),
    ("La aplicacion: 'Tienda de Perritos', un CRUD de productos de 3 capas.", 0),
    ("Frontend nginx (HTML/JS)", 1),
    ("Backend Node.js + Express (API REST)", 1),
    ("Base de datos MySQL 8", 1),
], size=19)
set_notes(s, "Contextualizar: esta es la tercera etapa. Ya habiamos contenedorizado; ahora el reto "
    "es orquestar en un cluster que escale y se recupere solo, y que el despliegue sea automatico. "
    "La app es sencilla a proposito (CRUD de productos) para concentrarnos en el DevOps.")

# ───────────────────────── 4. ARQUITECTURA (diagrama) ─────────────────────────
s = add_slide(); header_band(s, "Arquitectura de la solucion", "ECS Fargate + ALB + NLB")
# Internet
box(s, Inches(5.4), Inches(1.35), Inches(2.5), Inches(0.55), "Internet", DARK, size=14)
# ALB
box(s, Inches(4.3), Inches(2.15), Inches(4.7), Inches(0.95),
    "ALB publico  :80\n/ -> frontend     /api/* -> backend", BLUE, size=13)
# Frontend / Backend
box(s, Inches(2.0), Inches(3.5), Inches(3.2), Inches(1.0),
    "FRONTEND\nnginx :80  (2 tasks)", RGBColor(0x2563EB >> 16 & 255, 0x2563EB >> 8 & 255, 0x2563EB & 255), size=13)
box(s, Inches(7.9), Inches(3.5), Inches(3.4), Inches(1.0),
    "BACKEND\nnode :3001  (2 tasks)\nautoscaling CPU 50% (2->4)", RGBColor(0x25, 0x63, 0xEB), size=12)
# NLB
box(s, Inches(7.9), Inches(4.95), Inches(3.4), Inches(0.8), "NLB interno  TCP :3306", GRAY, size=13)
# DB
box(s, Inches(7.9), Inches(6.05), Inches(3.4), Inches(0.85),
    "MySQL DB  :3306  (1 task)\npassword desde SSM", DARK, size=12)
# flechas
arrow(s, Inches(6.65), Inches(1.9), Inches(6.65), Inches(2.15))
arrow(s, Inches(5.0), Inches(3.1), Inches(3.6), Inches(3.5))    # ALB -> front
arrow(s, Inches(8.2), Inches(3.1), Inches(9.6), Inches(3.5))    # ALB -> back
arrow(s, Inches(9.6), Inches(4.5), Inches(9.6), Inches(4.95))   # back -> NLB
arrow(s, Inches(9.6), Inches(5.75), Inches(9.6), Inches(6.05))  # NLB -> DB
_, tf = textbox(s, Inches(0.6), Inches(5.2), Inches(6.7), Inches(1.9))
for txt in ["Region us-east-1  -  Cuenta 166658674120",
            "Cluster: tienda-cluster (Fargate)",
            "Imagenes en Amazon ECR",
            "Logs en CloudWatch  -  Secret en SSM"]:
    p = tf.add_paragraph(); r = p.add_run(); r.text = "•  " + txt; style_run(r, 14, GRAY)
set_notes(s, "Esta es la foto completa. El navegador entra por el ALB publico. El ALB enruta por "
    "ruta: '/' al frontend y '/api/*' al backend (asi el front no necesita la IP del back). El "
    "backend habla con MySQL a traves de un NLB interno (TCP 3306). Todo corre en Fargate, las "
    "imagenes vienen de ECR, los logs van a CloudWatch y la clave de la BD vive en SSM. "
    "Insistir: NO usamos DNS interno de servicios porque el lab lo deniega; lo resolvimos con balanceadores.")

# ───────────────────────── 5. POR QUE FARGATE ─────────────────────────
s = add_slide(); header_band(s, "Por que ECS Fargate", "Decisiones tecnicas")
bullets(s, [
    ("ECS Fargate en vez de EKS:", 0),
    ("En el Learner Lab no se pueden crear roles IAM ni usar eksctl de forma fiable.", 1),
    ("Fargate es serverless: no administramos nodos EC2. Reutiliza el rol LabRole.", 1),
    ("Comunicacion por balanceadores en vez de Cloud Map / Service Connect:", 0),
    ("El lab deniega la accion servicediscovery:*.", 1),
    ("Front -> Back: ALB publico por ruta (/api/*).", 1),
    ("Back -> DB: NLB interno (capa 4, TCP 3306) con DNS estable.", 1),
    ("Resultado: arquitectura simple, escalable y 100% dentro de lo permitido por el lab.", 0),
], size=19)
set_notes(s, "Aca esta el 'por que' que mas preguntan. Dos decisiones clave nacieron de las "
    "restricciones del Learner Lab: (1) Fargate porque no podemos administrar roles/nodos como "
    "pide EKS, y (2) balanceadores en vez de DNS de servicio porque service discovery esta "
    "denegado. Son decisiones de ingenieria adaptadas al entorno real, no por capricho.")

# ───────────────────────── 6. CLUSTER (IE1) ─────────────────────────
s = add_slide(); header_band(s, "Configuracion del cluster", "Redes, Security Groups e IAM (IE1)")
bullets(s, [
    ("Cluster tienda-cluster, modo Fargate, red awsvpc (cada task con su propia IP).", 0),
    ("VPC vpc-00719b133c42a925b.", 0),
    ("Security Groups:", 0),
    ("tienda-alb-sg: entrante 80 desde Internet.", 1),
    ("tienda-ecs-sg: 80/3001 solo desde el ALB; 3306 solo dentro de la VPC.", 1),
    ("Roles IAM: LabRole como execution role (pull de ECR, leer SSM, escribir logs) y task role.", 0),
    ("En produccion real se separarian en roles distintos con privilegio minimo.", 1),
], size=19)
_, tf = textbox(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.7))
r = tf.paragraphs[0].add_run(); r.text = "[INSERTAR CAPTURA: consola ECS con los 3 servicios ACTIVE / reglas de los Security Groups]"
style_run(r, 13, ORANGE, italic=True)
set_notes(s, "El cluster usa awsvpc, o sea cada task tiene su propia interfaz de red e IP privada. "
    "La seguridad se controla con dos SG: el del ALB solo abre el 80 a Internet, y el de las tasks "
    "solo acepta trafico desde el ALB (80/3001) y el 3306 dentro de la VPC, asi la BD nunca queda "
    "expuesta. Sobre IAM: usamos LabRole para todo porque el lab no deja crear roles; aclarar que "
    "en produccion se separarian. Aca va una captura de la consola.")

# ───────────────────────── 7. DESPLIEGUE (IE2) ─────────────────────────
s = add_slide(); header_band(s, "Despliegue de Frontend y Backend", "Task Definitions desde ECR (IE2)")
bullets(s, [
    ("Cada servicio se define con una Task Definition de Fargate (carpeta ecs/).", 0),
    ("Imagenes desde Amazon ECR: tienda-frontend, tienda-backend, tienda-db.", 0),
    ("Variables de entorno del backend: DB_HOST (NLB), DB_USER, DB_NAME, DB_PORT.", 0),
    ("DB_PASSWORD inyectada como secret desde SSM (no en texto plano).", 0),
    ("Balanceo: target groups tg-frontend y tg-backend (ALB), tg-db (NLB).", 0),
    ("Acceso publico: http://tienda-alb-93590141.us-east-1.elb.amazonaws.com", 0),
], size=19)
_, tf = textbox(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.7))
r = tf.paragraphs[0].add_run(); r.text = "[INSERTAR CAPTURA: la app abierta en el navegador con la tabla de productos]"
style_run(r, 13, ORANGE, italic=True)
set_notes(s, "Cada contenedor se describe en una task definition: imagen de ECR, CPU/memoria, "
    "puertos, variables y secrets. El frontend escucha en 80, el backend en 3001. La contrasena "
    "de la BD entra como secret desde SSM. El acceso publico es la URL del ALB. Mostrar la app "
    "funcionando (captura o demo en vivo).")

# ───────────────────────── 8. COMUNICACION ─────────────────────────
s = add_slide(); header_band(s, "Comunicacion entre servicios")
box(s, Inches(0.8), Inches(2.4), Inches(2.6), Inches(1.0), "Navegador\n/api/productos", GRAY, size=13)
box(s, Inches(4.1), Inches(2.4), Inches(2.6), Inches(1.0), "ALB\nregla /api/* ", BLUE, size=13)
box(s, Inches(7.4), Inches(2.4), Inches(2.6), Inches(1.0), "Backend\n:3001", RGBColor(0x25,0x63,0xEB), size=13)
box(s, Inches(7.4), Inches(4.1), Inches(2.6), Inches(0.95), "NLB :3306", GRAY, size=13)
box(s, Inches(7.4), Inches(5.5), Inches(2.6), Inches(0.95), "MySQL", DARK, size=13)
arrow(s, Inches(3.4), Inches(2.9), Inches(4.1), Inches(2.9))
arrow(s, Inches(6.7), Inches(2.9), Inches(7.4), Inches(2.9))
arrow(s, Inches(8.7), Inches(3.4), Inches(8.7), Inches(4.1))
arrow(s, Inches(8.7), Inches(5.05), Inches(8.7), Inches(5.5))
bullets(s, [
    ("Front -> Back: el ALB enruta /api/* al backend. Sin DNS interno.", 0),
    ("Back -> DB: el backend usa el DNS del NLB como DB_HOST.", 0),
], t=Inches(6.0), size=17)
set_notes(s, "Detalle de como viaja una peticion: el navegador pide /api/productos, el ALB ve que "
    "calza con la regla /api/* y la manda al backend, el backend consulta MySQL a traves del NLB. "
    "Asi resolvimos el descubrimiento de servicios sin DNS interno, que estaba denegado.")

# ───────────────────────── 9. AUTOSCALING (IE3) ─────────────────────────
s = add_slide(); header_band(s, "Autoscaling", "Application Auto Scaling - Target Tracking (IE3)")
bullets(s, [
    ("Politica Target Tracking sobre CPU promedio del servicio.", 0),
    ("Objetivo: 50% de CPU.  Minimo 2 tasks, maximo 4.", 0),
    ("Por que 50%: deja margen para picos mientras se levantan tasks nuevas, sin sobre-aprovisionar.", 0),
    ("Cooldowns de 60 s para evitar oscilaciones (escalar/desescalar en cadena).", 0),
    ("Como se evidencia: generar carga (hey/ab) y ver en CloudWatch la CPU subir y ECS crear tasks.", 0),
], size=19)
_, tf = textbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.9))
r = tf.paragraphs[0].add_run()
r.text = "[INSERTAR CAPTURA: politica de autoscaling + grafico de CPU y recuento de tasks bajo carga]"
style_run(r, 13, ORANGE, italic=True)
set_notes(s, "El autoscaling mantiene la CPU promedio cerca del 50%: si sube, agrega tasks (hasta "
    "4); si baja, las quita (hasta 2). Elegimos 50% por margen de seguridad ante picos. Los "
    "cooldowns evitan que se ponga a escalar y desescalar nervioso. Para demostrarlo se hace una "
    "prueba de carga y se ve en CloudWatch como reacciona.")

# ───────────────────────── 10. PIPELINE (IE9) ─────────────────────────
s = add_slide(); header_band(s, "Pipeline CI/CD", "GitHub Actions: build -> push -> deploy (IE9)")
steps = ["push a\nmain", "Build\nimagenes", "Push a\nECR", "Render\ntask def", "Deploy\nECS"]
colors = [GRAY, BLUE, BLUE, RGBColor(0x25,0x63,0xEB), GREEN]
x = Inches(0.55)
for i, (txt, col) in enumerate(zip(steps, colors)):
    box(s, x, Inches(2.3), Inches(2.05), Inches(1.1), txt, col, size=14)
    if i < len(steps) - 1:
        arrow(s, x + Inches(2.05), Inches(2.85), x + Inches(2.45), Inches(2.85))
    x += Inches(2.5)
bullets(s, [
    ("Se dispara en cada push a main (y manualmente con workflow_dispatch).", 0),
    ("Tag de imagen = primeros 7 del commit SHA (trazabilidad).", 0),
    ("wait-for-service-stability: el deploy no es 'exito' hasta que ECS estabiliza las tasks.", 0),
], t=Inches(4.0), size=18)
set_notes(s, "El pipeline es el corazon del DevOps de la entrega. En cada push a main: hace build "
    "de las imagenes, las sube a ECR etiquetadas con el SHA del commit, renderiza la task "
    "definition con la imagen nueva y despliega en ECS esperando estabilidad. Si las tasks no "
    "quedan sanas, el paso falla. Aqui es donde, si el lab esta arriba, hacemos la DEMO en vivo: "
    "un commit pequeno y mostramos el Action corriendo.")

# ───────────────────────── 11. DEMO + METRICAS (IE6) ─────────────────────────
s = add_slide(); header_band(s, "Demo del pipeline y metricas", "Tiempos y ejecuciones reales (IE6)")
bullets(s, [
    ("Demo en vivo: hacer un commit y mostrar el run en GitHub Actions desplegando a ECS.", 0),
    ("4 ejecuciones registradas el 18-06-2026:", 0),
    ("Run #2 (exito): ~10 min   |   Run #4 (exito): ~16 min", 1),
    ("Run #1 (falla): secrets de AWS no cargados -> 'Could not load credentials'.", 1),
    ("Run #3 (falla): rollback automatico por el circuit breaker de ECS.", 1),
    ("El tiempo se va casi todo en la espera de estabilizacion del servicio.", 0),
    ("Las fallas fueron de configuracion / estabilizacion, NO del codigo (vistas en logs).", 0),
], size=18)
_, tf = textbox(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.6))
r = tf.paragraphs[0].add_run(); r.text = "[INSERTAR CAPTURA: pestaña Actions con el run en verde]"
style_run(r, 13, ORANGE, italic=True)
set_notes(s, "Mostrar metricas reales del pipeline: tarda 10-16 min, casi todo en estabilizar el "
    "servicio. Y contar la historia de las fallas porque queda muy bien en la defensa: la #1 fue "
    "por no haber cargado los secrets todavia, la #3 fue un rollback automatico del circuit "
    "breaker cuando las tasks no estabilizaron. Las diagnosticamos leyendo los logs de Actions. "
    "Esto demuestra que sabemos operar y depurar el pipeline, no solo armarlo.")

# ───────────────────────── 12. SECRETS (IE5) ─────────────────────────
s = add_slide(); header_band(s, "Gestion de secrets y credenciales", "(IE5)")
bullets(s, [
    ("Contrasena de la BD en AWS SSM Parameter Store (SecureString): /tienda/db_password.", 0),
    ("La task definition la referencia como secret -> el contenedor la recibe en runtime.", 1),
    ("No queda en la imagen, ni en el codigo, ni en el repositorio.", 1),
    ("Credenciales del pipeline como GitHub Actions Secrets (no en el codigo).", 0),
    (".gitignore excluye .env, *.pem y archivos de credenciales.", 0),
    ("Atencion: las credenciales del Learner Lab caducan; hay que refrescar los secrets del repo.", 0),
], size=19)
set_notes(s, "Punto de seguridad. La clave de la BD nunca se escribe en texto plano: vive en SSM "
    "como SecureString y se inyecta al contenedor solo en ejecucion. Las llaves de AWS para el "
    "pipeline son secrets de GitHub. Y el .gitignore evita subir credenciales por error. Mencionar "
    "que el session token del lab caduca y por eso hay que actualizar los secrets cada cierto tiempo.")

# ───────────────────────── 13. LOGS / VALIDACION (IE7) ─────────────────────────
s = add_slide(); header_band(s, "Logs, validacion y autorecuperacion", "(IE6 / IE7)")
bullets(s, [
    ("Logs: driver awslogs -> CloudWatch (/ecs/tienda-frontend, -backend, -db).", 0),
    ("Validacion funcional:", 0),
    ("Frontend accesible por la URL del ALB.", 1),
    ("/api/health responde ok; CRUD de productos opera contra MySQL.", 1),
    ("Comunicacion Front -> Back -> BD verificada.", 1),
    ("Autorecuperacion: si una task muere, ECS la reemplaza para mantener el desiredCount.", 0),
    ("Circuit breaker: ante un deploy defectuoso, ECS revierte solo a la version estable.", 0),
], size=18)
_, tf = textbox(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.6))
r = tf.paragraphs[0].add_run(); r.text = "[INSERTAR CAPTURA: CloudWatch logs + evento de task reemplazada]"
style_run(r, 13, ORANGE, italic=True)
set_notes(s, "Todos los contenedores mandan logs a CloudWatch, lo que nos permite depurar. Para "
    "validar mostramos la app funcionando, el health check y el CRUD. La tolerancia a fallos se ve "
    "en dos niveles: ECS repone tasks caidas, y el circuit breaker revierte despliegues malos solo. "
    "Si el lab esta arriba, se puede matar una task en vivo y mostrar como vuelve.")

# ───────────────────────── 14. PROBLEMAS (IE10) ─────────────────────────
s = add_slide(); header_band(s, "Problemas encontrados y soluciones", "Analisis critico (IE10)")
bullets(s, [
    ("servicediscovery denegado -> comunicacion con ALB (front-back) y NLB (back-db).", 0),
    ("No se pueden crear roles IAM / usar eksctl -> ECS Fargate con LabRole.", 0),
    ("Pipeline 'Could not load credentials' -> cargar los 3 secrets del lab en el repo.", 0),
    ("Rollback por circuit breaker -> relanzar deploy (quedo estable).", 0),
    ("docker login 400 en PowerShell -> usar docker login -u AWS -p <token>.", 0),
    ("ExpiredToken -> refrescar credenciales del Learner Lab.", 0),
], size=19)
set_notes(s, "Esta lamina es clave para IE10 (defensa). Cada fila es un problema real y como lo "
    "resolvimos. La narrativa: el entorno (Learner Lab) nos impuso restricciones y adaptamos la "
    "arquitectura en consecuencia; y cuando el pipeline fallo, diagnosticamos por logs y "
    "corregimos. Demuestra criterio tecnico, que es lo que evaluan en la defensa.")

# ───────────────────────── 15. CONCLUSIONES ─────────────────────────
s = add_slide(); header_band(s, "Conclusiones y proyeccion productiva")
bullets(s, [
    ("La app queda orquestada en ECS Fargate: escalable, tolerante a fallos y con CI/CD.", 0),
    ("El despliegue es automatico en cada push; las fallas son diagnosticables por logs.", 0),
    ("Proyeccion hacia produccion real para Innovatech Chile:", 0),
    ("Roles IAM separados con privilegio minimo (en vez de un LabRole unico).", 1),
    ("Base de datos en Amazon RDS gestionada (backups, multi-AZ) en vez de contenedor.", 1),
    ("HTTPS en el ALB (ACM) + dominio propio en Route 53.", 1),
    ("OIDC entre GitHub y AWS (sin claves estaticas que caducan).", 1),
    ("Alarmas de CloudWatch y entorno de staging previo a produccion.", 1),
], size=18)
set_notes(s, "Cerrar con el balance: cumplimos los objetivos (escalabilidad, tolerancia a fallos, "
    "automatizacion). Y mostrar vision: que haria falta para llevar esto a produccion de verdad. "
    "Esto suele sumar puntos porque demuestra que entendemos las limitaciones del entorno academico.")

# ───────────────────────── 16. CIERRE ─────────────────────────
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = BLUE; bg.line.fill.background(); bg.shadow.inherit = False
_, tf = textbox(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.2), anchor=MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "Gracias"
style_run(r, 54, WHITE, bold=True)
p = tf.add_paragraph(); r = p.add_run()
r.text = "Preguntas?   -   github.com/jeffplop/tienda-perritos"
style_run(r, 22, RGBColor(0xDB, 0xEA, 0xFE))
set_notes(s, "Agradecer y abrir a preguntas. Tener a mano el repo, la consola de AWS y la pestana "
    "Actions para responder mostrando evidencia. Recordar el guion de preguntas de defensa.")

out = "/home/user/tienda-perritos/docs/EP3-Presentacion-Tienda-Perritos.pptx"
prs.save(out)
print("OK ->", out, "  laminas:", len(prs.slides._sldIdLst))
